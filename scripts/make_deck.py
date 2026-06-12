#!/usr/bin/env python3
"""Generate 'THE VILLA REPORT' PowerPoint from the live database.

Re-run any time (e.g. after the 6pm pull) to refresh every number and quote:
    .venv/bin/python scripts/make_deck.py
Output: outputs/love_island_report.pptx
"""

from __future__ import annotations

import argparse
import math
import re
import shutil
import sys
from datetime import datetime, timezone
from pathlib import Path
from zoneinfo import ZoneInfo

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

import matplotlib
matplotlib.use("Agg")
import matplotlib.dates as mdates
import matplotlib.pyplot as plt
import pandas as pd

from loveisland.store import db

CHARTS = Path(__file__).resolve().parent.parent / "outputs" / "_charts"

# ── Palette (Love Island after-dark) ──────────────────────────────────
BG    = RGBColor(0x11, 0x0D, 0x22)
CARD  = RGBColor(0x1E, 0x16, 0x34)
CARD2 = RGBColor(0x2A, 0x1F, 0x47)
PINK  = RGBColor(0xFF, 0x2E, 0x88)
GOLD  = RGBColor(0xFF, 0xC9, 0x3C)
WHITE = RGBColor(0xF5, 0xF2, 0xFA)
MUTE  = RGBColor(0x9E, 0x8F, 0xB8)
GREEN = RGBColor(0x36, 0xE2, 0x7B)
RED   = RGBColor(0xFF, 0x5C, 0x5C)
HEAD, BODY = "Arial Black", "Arial"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)


# ── Data ───────────────────────────────────────────────────────────────

def fetch():
    c = db.connect()
    q = lambda s: c.execute(s).fetchall()
    d = {}
    d["total"] = c.execute("SELECT COUNT(*) FROM items WHERE sentiment_label IS NOT NULL").fetchone()[0]
    d["avg"] = c.execute("SELECT AVG(sentiment_score) FROM items WHERE sentiment_score IS NOT NULL").fetchone()[0] or 0.0
    d["sources"] = q("SELECT source, COUNT(*) n FROM items GROUP BY source ORDER BY n DESC")
    d["span"] = c.execute("SELECT MIN(created_at), MAX(created_at) FROM items").fetchone()
    d["pos"] = c.execute("SELECT AVG(sentiment_score>=0.15) FROM items WHERE sentiment_score IS NOT NULL").fetchone()[0] or 0
    d["neg"] = c.execute("SELECT AVG(sentiment_score<=-0.15) FROM items WHERE sentiment_score IS NOT NULL").fetchone()[0] or 0
    d["discussed"] = q("SELECT entity, COUNT(*) n FROM aspects WHERE entity_type='contestant' GROUP BY entity ORDER BY n DESC LIMIT 6")
    d["faves"] = q("SELECT entity, COUNT(*) n, AVG(sentiment_score) a FROM aspects WHERE entity_type='contestant' GROUP BY entity HAVING n>=5 ORDER BY a DESC LIMIT 5")
    d["villains"] = q("SELECT entity, COUNT(*) n, AVG(sentiment_score) a FROM aspects WHERE entity_type='contestant' GROUP BY entity HAVING n>=5 ORDER BY a ASC LIMIT 5")
    d["report"] = q("SELECT entity, COUNT(*) n, AVG(sentiment_score) a FROM aspects WHERE entity_type='contestant' GROUP BY entity HAVING n>=5 ORDER BY a DESC LIMIT 15")
    d["all_contestants"] = q("SELECT entity, COUNT(*) n, AVG(sentiment_score) a FROM aspects WHERE entity_type='contestant' GROUP BY entity HAVING n>=5 ORDER BY n DESC")
    d["couples"] = q("SELECT entity, COUNT(*) n, AVG(sentiment_score) a FROM aspects WHERE entity_type='couple' GROUP BY entity HAVING n>=3 ORDER BY a DESC")
    d["burns"] = q("SELECT text, source FROM items WHERE funny>=0.5 AND sentiment_score<=-0.2 AND LENGTH(text) BETWEEN 25 AND 210 ORDER BY funny DESC, like_count DESC LIMIT 6")
    d["funny"] = q("SELECT text, source FROM items WHERE funny>=0.5 AND sentiment_score>-0.2 AND LENGTH(text) BETWEEN 25 AND 210 ORDER BY funny DESC, like_count DESC LIMIT 6")
    return d


def clean(t, n=300):
    t = re.sub(r"\s+", " ", (t or "").strip())
    return (t[:n].rstrip() + "…") if len(t) > n else t


def fit_size(txt, w, h, lo, hi):
    """Largest font (pt) in [lo, hi] that fits `txt` in a w x h inch area."""
    pt = math.sqrt(7000 * w * h / max(1, len(txt)))
    return round(max(lo, min(hi, pt)), 1)


def grade(a):
    if a >= 0.22: return "A", GREEN
    if a >= 0.08: return "B", GREEN
    if a >= -0.08: return "C", GOLD
    if a >= -0.22: return "D", RED
    return "F", RED


# ── Slide helpers ───────────────────────────────────────────────────────

def slide(prs, bg=BG):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background(); r.shadow.inherit = False
    return s


def text(s, l, t, w, h, content, size, color, *, font=BODY, bold=True, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, spacing=1.0):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align; p.line_spacing = spacing
    r = p.add_run(); r.text = content
    r.font.size = Pt(size); r.font.bold = bold; r.font.name = font; r.font.color.rgb = color
    return tb


def card(s, l, t, w, h, fill=CARD):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def bar(s, l, t, w, h, frac, color):
    tr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    tr.fill.solid(); tr.fill.fore_color.rgb = CARD2; tr.line.fill.background(); tr.shadow.inherit = False
    fw = max(0.12, w * max(0.04, min(1.0, frac)))
    fl = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(fw), Inches(h))
    fl.fill.solid(); fl.fill.fore_color.rgb = color; fl.line.fill.background(); fl.shadow.inherit = False


def title(s, kicker, ttl, ttl_color=WHITE):
    text(s, 0.62, 0.42, 12.1, 0.4, kicker.upper(), 15, GOLD, font=BODY)
    text(s, 0.6, 0.74, 12.1, 1.0, ttl.upper(), 38, ttl_color, font=HEAD)


def ranking(s, rows, top, color_fn, denom=0.4):
    for i, (name, n, a) in enumerate(rows):
        y = top + i * 0.86
        text(s, 0.7, y, 3.6, 0.5, name, 25, WHITE, font=HEAD, anchor=MSO_ANCHOR.MIDDLE)
        text(s, 0.7, y + 0.5, 3.6, 0.3, f"{n:,} mentions", 11, MUTE)
        col = color_fn(a)
        bar(s, 4.5, y + 0.12, 6.7, 0.4, abs(a) / denom, col)
        text(s, 11.35, y, 1.7, 0.6, f"{a:+.2f}", 22, col, font=HEAD, anchor=MSO_ANCHOR.MIDDLE)


def quotes_grid(s, rows):
    """6 quote cards, 2 columns x 3 rows; full quotes, auto-sized to fit."""
    pos = [(0.6, 1.85), (6.85, 1.85), (0.6, 3.72), (6.85, 3.72), (0.6, 5.59), (6.85, 5.59)]
    for (l, t), (txt, src) in zip(pos, rows):
        card(s, l, t, 5.88, 1.7)
        text(s, l + 0.28, t + 0.1, 0.6, 0.4, "“", 28, PINK, font=HEAD)
        q = clean(txt)
        text(s, l + 0.3, t + 0.42, 5.32, 0.96, q, fit_size(q, 5.0, 0.95, 10.5, 14), WHITE, bold=False, spacing=1.03)
        text(s, l + 0.3, t + 1.4, 5.32, 0.26, f"via {src}", 10.5, GOLD)


# ── Trend charts + deep dives ───────────────────────────────────────────

def _series(rows):
    if not rows:
        return None
    df = pd.DataFrame(rows, columns=["t", "s"])
    df["t"] = pd.to_datetime(df["t"], utc=True, errors="coerce")
    ser = df.dropna(subset=["t"]).set_index("t").resample("1D")["s"].mean().dropna()
    return ser if len(ser) >= 2 else None


def overall_series():
    rows = db.connect().execute(
        "SELECT created_at, sentiment_score FROM items WHERE sentiment_score IS NOT NULL").fetchall()
    return _series(rows)


def entity_series(entity):
    rows = db.connect().execute(
        "SELECT i.created_at, a.sentiment_score FROM aspects a JOIN items i ON i.id=a.item_id "
        "WHERE a.entity=?", (entity,)).fetchall()
    return _series(rows)


def trend_png(ser, color_hex, name, h_in=2.5):
    CHARTS.mkdir(parents=True, exist_ok=True)
    path = CHARTS / (re.sub(r"[^A-Za-z0-9]+", "_", name) + ".png")
    fig, ax = plt.subplots(figsize=(11.9, h_in), dpi=150)
    fig.patch.set_alpha(0); ax.set_facecolor("none")
    ax.axhline(0, color="#5A4F73", lw=1.2)
    ax.fill_between(ser.index, ser.values, 0, color="#" + color_hex, alpha=0.13)
    ax.plot(ser.index, ser.values, color="#" + color_hex, lw=3.4, marker="o", ms=6,
            mfc="#" + color_hex, mec="white", mew=0.6)
    ax.set_ylim(-1, 1); ax.set_yticks([-1, -0.5, 0, 0.5, 1])
    loc = mdates.AutoDateLocator()
    ax.xaxis.set_major_locator(loc)
    ax.xaxis.set_major_formatter(mdates.ConciseDateFormatter(loc))
    ax.tick_params(colors="#9E8FB8", labelsize=10)
    for sp in ("top", "right"): ax.spines[sp].set_visible(False)
    for sp in ("left", "bottom"): ax.spines[sp].set_color("#5A4F73")
    ax.grid(axis="y", color="#2A1F47", lw=0.9)
    fig.tight_layout(pad=0.5)
    fig.savefig(path, transparent=True); plt.close(fig)
    return str(path)


def deep_dive(prs, name, kind, n, a):
    s = slide(prs)
    gl, gc = grade(a)
    text(s, 0.62, 0.42, 10.0, 0.4, f"APPENDIX · {kind} DEEP DIVE", 13, GOLD)
    text(s, 0.6, 0.76, 10.0, 0.9, name.upper(), 33, WHITE, font=HEAD)
    text(s, 0.62, 1.56, 10.0, 0.4, f"{n:,} mentions   ·   avg sentiment {a:+.2f}", 14, MUTE, bold=False)
    card(s, 11.35, 0.5, 1.45, 1.45, CARD)
    text(s, 11.35, 0.5, 1.45, 1.45, gl, 58, gc, font=HEAD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    text(s, 0.62, 2.02, 11.9, 0.32, "SENTIMENT OVER TIME", 13, MUTE)
    ser = entity_series(name)
    if ser is not None:
        png = trend_png(ser, "36E27B" if a >= 0 else "FF5C5C", name, 2.4)
        s.shapes.add_picture(png, Inches(0.7), Inches(2.36), Inches(11.9), Inches(2.5))
    else:
        card(s, 0.7, 2.4, 11.9, 2.4, CARD)
        text(s, 0.7, 3.35, 11.9, 0.6, "Not enough history yet for a trend line — fills in over the season.",
             16, MUTE, bold=False, align=PP_ALIGN.CENTER)

    c = db.connect()
    fit = ("SELECT i.text, i.source FROM aspects a JOIN items i ON i.id=a.item_id "
           "WHERE a.entity=? AND LENGTH(i.text) BETWEEN 25 AND 260 "
           "ORDER BY a.sentiment_score {d}, i.like_count DESC LIMIT 1")
    anyq = ("SELECT i.text, i.source FROM aspects a JOIN items i ON i.id=a.item_id "
            "WHERE a.entity=? ORDER BY a.sentiment_score {d}, i.like_count DESC LIMIT 1")
    pick = lambda dr: c.execute(fit.format(d=dr), (name,)).fetchone() or c.execute(anyq.format(d=dr), (name,)).fetchone()
    loved, dragged = pick("DESC"), pick("ASC")
    for l, head, hc, row in [(0.6, "💖 BEST TAKE", GREEN, loved), (6.85, "🔥 WORST TAKE", RED, dragged)]:
        card(s, l, 5.1, 5.88, 2.0, CARD)
        text(s, l + 0.3, 5.26, 5.3, 0.32, head, 12, hc)
        if row:
            q = clean(row[0])
            text(s, l + 0.3, 5.62, 5.3, 1.12, q, fit_size(q, 5.0, 1.12, 11, 15), WHITE, bold=False, spacing=1.04)
            text(s, l + 0.3, 6.78, 5.3, 0.26, f"via {row[1]}", 10.5, GOLD)
        else:
            text(s, l + 0.3, 5.7, 5.3, 0.5, "—", 14, MUTE, bold=False)


# ── "What changed since last time" ──────────────────────────────────────

def fmt_baseline(iso):
    try:
        return datetime.fromisoformat(iso).astimezone(ZoneInfo("America/New_York")).strftime("%b %-d, %-I%p")
    except Exception:
        return "last time"


def change_metrics(baseline):
    c = db.connect()
    b = baseline

    def wins(et):
        new = {e: (n, a) for e, n, a in c.execute(
            "SELECT a.entity,COUNT(*),AVG(a.sentiment_score) FROM aspects a JOIN items i ON i.id=a.item_id "
            "WHERE i.created_at>? AND a.entity_type=? GROUP BY a.entity", (b, et))}
        old = {e: (n, a) for e, n, a in c.execute(
            "SELECT a.entity,COUNT(*),AVG(a.sentiment_score) FROM aspects a JOIN items i ON i.id=a.item_id "
            "WHERE i.created_at<=? AND a.entity_type=? GROUP BY a.entity", (b, et))}
        return new, old

    M = {}
    M["new_count"] = c.execute("SELECT COUNT(*) FROM items WHERE created_at>? AND sentiment_label IS NOT NULL", (b,)).fetchone()[0]
    M["avg_new"] = c.execute("SELECT AVG(sentiment_score) FROM items WHERE created_at>? AND sentiment_score IS NOT NULL", (b,)).fetchone()[0] or 0.0
    M["avg_old"] = c.execute("SELECT AVG(sentiment_score) FROM items WHERE created_at<=? AND sentiment_score IS NOT NULL", (b,)).fetchone()[0] or 0.0
    cn, co = wins("contestant")
    mov = [(e, co[e][1], an, an - co[e][1]) for e, (nn, an) in cn.items() if nn >= 8 and e in co and co[e][0] >= 8]
    M["risers"] = sorted([x for x in mov if x[3] > 0.03], key=lambda x: -x[3])[:4]
    M["fallers"] = sorted([x for x in mov if x[3] < -0.03], key=lambda x: x[3])[:4]
    M["buzz"] = sorted(cn.items(), key=lambda kv: -kv[1][0])[:1]
    pn, po = wins("couple")
    cm = [(e, po[e][1], an, an - po[e][1]) for e, (nn, an) in pn.items() if nn >= 4 and e in po and po[e][0] >= 4]
    M["couples"] = sorted(cm, key=lambda x: -abs(x[3]))[:3]
    M["fresh_burns"] = c.execute(
        "SELECT text,source FROM items WHERE created_at>? AND funny>=0.5 AND sentiment_score<=-0.2 "
        "AND LENGTH(text) BETWEEN 25 AND 210 ORDER BY funny DESC, like_count DESC LIMIT 2", (b,)).fetchall()
    return M


def _mover_rows(s, l, t, head, col, rows):
    card(s, l, t, 5.9, 3.45)
    text(s, l + 0.3, t + 0.22, 5.3, 0.4, head, 18, col, font=HEAD)
    if not rows:
        text(s, l + 0.3, t + 1.4, 5.3, 0.5, "No big swings this time.", 14, MUTE, bold=False)
        return
    for i, (e, ao, an, d) in enumerate(rows):
        y = t + 0.8 + i * 0.62
        text(s, l + 0.3, y, 2.3, 0.5, e, 18, WHITE, font=HEAD, anchor=MSO_ANCHOR.MIDDLE)
        text(s, l + 2.55, y, 2.05, 0.5, f"{ao:+.2f} → {an:+.2f}", 14, MUTE, anchor=MSO_ANCHOR.MIDDLE)
        text(s, l + 4.6, y, 1.1, 0.5, f"{d:+.2f}", 17, col, font=HEAD, anchor=MSO_ANCHOR.MIDDLE)


def change_slides(prs, M, baseline):
    lbl = fmt_baseline(baseline)
    s = slide(prs)
    title(s, f"what's moved since {lbl}", "Since We Last Looked 🔄")
    delta = M["avg_new"] - M["avg_old"]
    arrow, col = ("▼", RED) if delta < -0.01 else (("▲", GREEN) if delta > 0.01 else ("→", GOLD))
    verdict = "the villa SOURED" if delta < -0.02 else ("the villa WARMED UP" if delta > 0.02 else "barely budged")
    card(s, 0.7, 1.82, 11.95, 1.38, CARD2)
    text(s, 1.1, 1.98, 7.4, 0.36, "THE VILLA'S MOOD", 14, MUTE)
    text(s, 1.1, 2.34, 7.4, 0.8, f"{M['avg_old']:+.2f}   →   {M['avg_new']:+.2f}  {arrow}", 35, col, font=HEAD)
    text(s, 8.7, 2.05, 3.8, 0.5, verdict, 18, WHITE, font=HEAD)
    text(s, 8.7, 2.58, 3.8, 0.5, f"{M['new_count']:,} new comments", 14, MUTE, bold=False)
    _mover_rows(s, 0.7, 3.45, "📈 RISING", GREEN, M["risers"])
    _mover_rows(s, 6.75, 3.45, "📉 FALLING", RED, M["fallers"])

    s = slide(prs)
    title(s, "who blew up and the freshest tea", "The Episode Effect 🎬")
    if M["buzz"]:
        e, (nn, _a) = M["buzz"][0]
        card(s, 0.7, 1.82, 11.95, 1.2, CARD2)
        text(s, 1.1, 1.98, 7.0, 0.35, "🗣️ MAIN CHARACTER OF THE NIGHT", 14, MUTE)
        text(s, 1.1, 2.32, 8.2, 0.7, e, 34, PINK, font=HEAD)
        text(s, 9.0, 2.02, 3.4, 0.5, f"{nn:,} new mentions", 16, WHITE, font=HEAD, anchor=MSO_ANCHOR.MIDDLE)
        text(s, 9.0, 2.52, 3.4, 0.45, "all eyes on them", 13, MUTE, bold=False)
    text(s, 0.7, 3.3, 11.9, 0.35, "💔 COUPLE SHIFTS", 14, MUTE)
    if M["couples"]:
        for i, (e, ao, an, d) in enumerate(M["couples"]):
            y = 3.72 + i * 0.48
            c2 = GREEN if d >= 0 else RED
            text(s, 0.9, y, 4.8, 0.4, f"{'▲' if d >= 0 else '▼'} {e}", 16, c2, font=HEAD, anchor=MSO_ANCHOR.MIDDLE)
            text(s, 6.0, y, 2.6, 0.4, f"{ao:+.2f} → {an:+.2f}", 14, MUTE, anchor=MSO_ANCHOR.MIDDLE)
            text(s, 8.8, y, 1.4, 0.4, f"{d:+.2f}", 15, c2, font=HEAD, anchor=MSO_ANCHOR.MIDDLE)
    else:
        text(s, 0.9, 3.72, 11.0, 0.4, "No notable couple swings since last time.", 14, MUTE, bold=False)
    text(s, 0.7, 5.25, 11.9, 0.35, "🔥 FRESH OFF THE EPISODE", 14, MUTE)
    for (l, t), (txt, src) in zip([(0.6, 5.62), (6.85, 5.62)], M["fresh_burns"]):
        card(s, l, t, 5.88, 1.5)
        q = clean(txt)
        text(s, l + 0.3, t + 0.16, 5.3, 0.92, q, fit_size(q, 5.0, 0.9, 10.5, 13.5), WHITE, bold=False, spacing=1.02)
        text(s, l + 0.3, t + 1.18, 5.3, 0.26, f"via {src}", 10, GOLD)


# ── Build ────────────────────────────────────────────────────────────────

def build(d, baseline=None):
    prs = Presentation(); prs.slide_width = EMU_W; prs.slide_height = EMU_H
    span_hi = (d["span"][1] or "")[:10]
    best = d["couples"][0]; worst_c = d["couples"][-1]
    worst = d["villains"][0][0]; fave = d["faves"][0][0]
    drama = next((e for e, n, a in d["villains"]), worst)

    # 1 — Title
    s = slide(prs)
    card(s, 0, 0, 13.333, 0.16, PINK)
    text(s, 0.9, 1.7, 11.5, 0.5, "LOVE ISLAND USA · SEASON 8", 18, GOLD, align=PP_ALIGN.CENTER)
    text(s, 0.5, 2.35, 12.3, 1.1, "THE VILLA REPORT", 52, WHITE, font=HEAD, align=PP_ALIGN.CENTER)
    text(s, 1.2, 3.7, 10.9, 0.7, "What the internet REALLY thinks before tonight's episode 🌴🍷",
         22, PINK, align=PP_ALIGN.CENTER)
    text(s, 1.2, 5.25, 10.9, 0.6, f"{d['total']:,} comments read across YouTube · Bluesky · Tumblr · News",
         16, WHITE, bold=False, align=PP_ALIGN.CENTER)
    text(s, 1.2, 5.72, 10.9, 0.5, "so you don't have to.", 14, MUTE, bold=False, align=PP_ALIGN.CENTER)

    # 2 — Vibe Check
    s = slide(prs)
    title(s, "the state of the villa", "The Vibe Check")
    days = (datetime.fromisoformat(d["span"][1]).date() - datetime.fromisoformat(d["span"][0]).date()).days
    for i, (big, lab) in enumerate([(f"{d['total']:,}", "comments analyzed"), (str(len(d["sources"])), "platforms"), (f"{days}d", "of season covered")]):
        x = 0.7 + i * 4.15
        card(s, x, 2.0, 3.8, 1.7)
        text(s, x, 2.2, 3.8, 0.95, big, 50, GOLD, font=HEAD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, x, 3.15, 3.8, 0.4, lab, 14, MUTE, align=PP_ALIGN.CENTER)
    mood = "DIVIDED 😬" if abs(d["avg"]) < 0.1 else ("ROUGH 🔥" if d["avg"] < 0 else "GLOWING ☀️")
    card(s, 0.7, 4.1, 11.95, 2.55, CARD2)
    text(s, 1.1, 4.4, 11.2, 0.5, "OVERALL MOOD", 15, MUTE)
    text(s, 1.1, 4.78, 11.2, 1.0, mood, 46, PINK, font=HEAD)
    text(s, 1.1, 5.85, 11.2, 0.7, f"{d['pos']*100:.0f}% feeling good  ·  {d['neg']*100:.0f}% absolutely seething  ·  the internet is FEASTING.", 17, WHITE, bold=False)

    # 3 & 4 — What changed since last time (only when we have a baseline)
    if baseline:
        cm = change_metrics(baseline)
        if cm["new_count"] > 0:
            change_slides(prs, cm, baseline)

    # The Mood Swing (overall sentiment over time)
    s = slide(prs)
    title(s, "how the villa's mood has swung all season", "The Mood Swing 📈")
    ser = overall_series()
    if ser is not None:
        png = trend_png(ser, "FF2E88", "_overall", 3.1)
        s.shapes.add_picture(png, Inches(0.7), Inches(2.1), Inches(11.9), Inches(3.25))
        cur, lo, hi = ser.iloc[-1], ser.min(), ser.max()
        vibe = "happy" if cur >= 0.1 else ("rough" if cur <= -0.1 else "split down the middle")
        text(s, 0.7, 5.65, 11.9, 0.9,
             f"Right now the villa is {vibe} ({cur:+.2f}). It's swung from {lo:+.2f} to {hi:+.2f} this season — a total rollercoaster.",
             16, WHITE, bold=False)
    else:
        text(s, 0.7, 3.3, 11.9, 0.6, "Not enough history yet — fills in as the season goes.", 18, MUTE, bold=False)

    # 4 — Main Characters
    s = slide(prs)
    title(s, "who the internet can't shut up about", "Main Characters")
    text(s, 0.62, 1.72, 12.1, 0.5, f"{d['discussed'][0][0]} is living rent-free in everyone's head.", 16, GOLD, bold=False)
    mx = max(n for _, n in d["discussed"]) or 1
    for i, (name, n) in enumerate(d["discussed"]):
        y = 2.35 + i * 0.78
        text(s, 0.7, y, 3.4, 0.5, name, 23, WHITE, font=HEAD, anchor=MSO_ANCHOR.MIDDLE)
        bar(s, 4.3, y + 0.1, 6.9, 0.38, n / mx, PINK)
        text(s, 11.35, y, 1.8, 0.5, f"{n:,}", 20, GOLD, font=HEAD, anchor=MSO_ANCHOR.MIDDLE)

    # 4 — Fan Favorites
    s = slide(prs)
    title(s, "the ones who can do no wrong", "Fan Favorites 💖", GREEN)
    text(s, 0.62, 1.72, 12.1, 0.5, f"{fave} is the people's princess.", 16, GOLD, bold=False)
    ranking(s, d["faves"], 2.4, lambda a: GREEN)

    # 5 — Getting Cooked
    s = slide(prs)
    title(s, "public enemies", "Getting Cooked 🔥", RED)
    text(s, 0.62, 1.72, 12.1, 0.5, f"{worst} is public enemy #1. The villa wants blood.", 16, GOLD, bold=False)
    ranking(s, d["villains"], 2.4, lambda a: RED)

    # 6 — Cast Report Card
    s = slide(prs)
    title(s, "every islander, graded by the fans", "The Report Card 📋")
    for i, (name, n, a) in enumerate(d["report"]):
        gl, gc = grade(a)
        col, row = i % 3, i // 3
        x, y = 0.6 + col * 4.12, 2.0 + row * 0.95
        card(s, x, y, 3.9, 0.82)
        text(s, x + 0.25, y + 0.1, 2.55, 0.38, name, 17, WHITE, font=HEAD, anchor=MSO_ANCHOR.MIDDLE)
        text(s, x + 0.25, y + 0.5, 2.55, 0.26, f"{n:,} mentions", 9.5, MUTE)
        text(s, x + 2.85, y, 0.9, 0.82, gl, 36, gc, font=HEAD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 7 — Couple Chemistry
    s = slide(prs)
    title(s, "endgame or endgame-over", "Couple Chemistry 💑")
    text(s, 0.62, 1.72, 12.1, 0.6, f"The people's verdict: {best[0]} 💚 is winning, {worst_c[0]} 💔 is doomed.", 15, GOLD, bold=False)
    cx, cw = 6.7, 4.7
    for i, (name, n, a) in enumerate(d["couples"]):
        y = 2.45 + i * 0.5
        col = GREEN if a >= 0 else RED
        text(s, 0.7, y - 0.02, 4.4, 0.4, name, 15, WHITE, anchor=MSO_ANCHOR.MIDDLE)
        frac = max(0.03, min(1.0, abs(a) / 0.45)) * cw
        if a >= 0:
            bar(s, cx, y + 0.05, frac, 0.28, 1.0, col)
        else:
            b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx - frac), Inches(y + 0.05), Inches(frac), Inches(0.28))
            b.fill.solid(); b.fill.fore_color.rgb = col; b.line.fill.background(); b.shadow.inherit = False
        text(s, cx + cw + 0.15, y - 0.02, 1.0, 0.4, f"{a:+.2f}", 13, col, font=HEAD, anchor=MSO_ANCHOR.MIDDLE)
    zl = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx - 0.01), Inches(2.4), Inches(0.02), Inches(len(d["couples"]) * 0.5 + 0.1))
    zl.fill.solid(); zl.fill.fore_color.rgb = MUTE; zl.line.fill.background(); zl.shadow.inherit = False

    # 8 — The Burn Book
    s = slide(prs)
    title(s, "the most savage takes we found", "The Burn Book 😈")
    quotes_grid(s, d["burns"])

    # 9 — Funniest Takes
    s = slide(prs)
    title(s, "comments that made us actually laugh", "Funniest Takes 😂")
    quotes_grid(s, d["funny"])

    # 10 — The Odds
    s = slide(prs)
    title(s, "what the data is betting on tonight", "The Odds 🎲")
    preds = [
        ("🎬", "Most likely to cause a scene", drama),
        ("🚪", "Most likely to get dumped (per the haters)", worst),
        ("💔", "Couple most likely to implode", worst_c[0]),
        ("💍", "Ship the people are manifesting", best[0]),
        ("⭐", "Dark horse winning everyone over", fave),
    ]
    for i, (ico, lab, who) in enumerate(preds):
        y = 1.95 + i * 0.98
        card(s, 0.7, y, 11.95, 0.84)
        text(s, 0.95, y, 0.9, 0.84, ico, 26, WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, 1.95, y, 7.0, 0.84, lab, 18, WHITE, bold=False, anchor=MSO_ANCHOR.MIDDLE)
        text(s, 9.0, y, 3.45, 0.84, who, 22, PINK, font=HEAD, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    # 11 — What To Watch
    s = slide(prs)
    title(s, "storylines the data says to watch", "What To Watch Tonight 📺")
    blocks = [
        ("🎯 The Melanie Pile-On", "She's the most-talked-about islander by a mile and the fans are NOT kind. Every move gets clipped."),
        ("💘 The Forbidden Ship", "Fans want Sol & Sincere over his actual couple Melanie & Sincere. The triangle's about to blow."),
        ("🐍 Gabriel, Menace at Large", "Kissing everyone, trusted by no one. One of the lowest-rated men in the villa."),
        ("👬 Zryce Nation", "The internet has decided Zach & Bryce are the real love story. Let them cook."),
    ]
    for i, (h, b) in enumerate(blocks):
        x = 0.7 + (i % 2) * 6.25; y = 2.05 + (i // 2) * 2.35
        card(s, x, y, 5.7, 2.05)
        text(s, x + 0.3, y + 0.25, 5.1, 0.6, h, 19, PINK, font=HEAD)
        text(s, x + 0.3, y + 0.95, 5.1, 1.0, b, 14, WHITE, bold=False, spacing=1.05)

    # 12 — The Verdict
    s = slide(prs)
    card(s, 0, 0, 13.333, 0.16, PINK)
    text(s, 0.7, 1.3, 11.9, 0.9, "THE VERDICT", 50, WHITE, font=HEAD, align=PP_ALIGN.CENTER)
    for i, (lab, val, col) in enumerate([("💖 FAN FAVORITE", fave, GREEN), ("🔥 MOST BOOED", worst, RED), ("💑 SHIP THE PEOPLE DEMAND", best[0], PINK)]):
        x = 0.7 + i * 4.15
        card(s, x, 2.7, 3.8, 1.9)
        text(s, x + 0.2, 2.95, 3.4, 0.4, lab, 13, MUTE, align=PP_ALIGN.CENTER)
        text(s, x + 0.2, 3.45, 3.4, 1.0, val, 27, col, font=HEAD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 1.0, 5.2, 11.3, 0.7, "Now stop reading and go watch. 🍷📺", 26, GOLD, font=HEAD, align=PP_ALIGN.CENTER)
    text(s, 1.0, 6.1, 11.3, 0.5, f"Data as of {span_hi} · {d['total']:,} comments analyzed · the people have spoken", 12, MUTE, bold=False, align=PP_ALIGN.CENTER)

    # ── APPENDIX: deep dive per contestant + couple ──────────────────────
    s = slide(prs)
    card(s, 0, 3.5, 13.333, 0.16, PINK)
    text(s, 0.7, 2.45, 11.9, 1.1, "APPENDIX", 62, PINK, font=HEAD, align=PP_ALIGN.CENTER)
    text(s, 0.7, 3.8, 11.9, 0.6, "every contestant & couple, deep-dived 🔬", 20, MUTE, align=PP_ALIGN.CENTER, bold=False)
    for name, n, a in d["all_contestants"]:
        deep_dive(prs, name, "CONTESTANT", n, a)
    for name, n, a in d["couples"]:
        deep_dive(prs, name, "COUPLE", n, a)

    return prs


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--preview", action="store_true", help="don't update the 'last looked' snapshot")
    ap.add_argument("--since", help="override the comparison baseline (ISO timestamp)")
    args = ap.parse_args()

    data = fetch()
    baseline = args.since or db.get_meta("deck_last_run")
    out = Path(__file__).resolve().parent.parent / "outputs" / "love_island_report.pptx"
    out.parent.mkdir(exist_ok=True)
    prs = build(data, baseline)
    prs.save(out)
    note = ""
    if not args.preview:  # production run: bookmark "now" and archive a versioned copy
        db.set_meta("deck_last_run", datetime.now(timezone.utc).isoformat())
        archive_dir = out.parent / "decks"
        archive_dir.mkdir(parents=True, exist_ok=True)
        stamp = datetime.now().astimezone().strftime("%Y-%m-%d_%H%M")
        archived = archive_dir / f"villa_report_{stamp}.pptx"
        shutil.copyfile(out, archived)
        note = f"\n   archived version → {archived}"
    tag = f" · vs {fmt_baseline(baseline)}" if baseline else ""
    print(f"✅ Saved {out}  ({data['total']:,} comments, {len(prs.slides)} slides{tag}){note}")
